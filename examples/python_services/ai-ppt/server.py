"""ai-ppt gRPC server — generate PowerPoint via LLM"""
import asyncio, os
from typing import AsyncIterator
import grpc
from grpc.aio import server as grpc_server
from pptx import Presentation
import litellm
import ppt_pb2, ppt_pb2_grpc

class PptService(ppt_pb2_grpc.PptServiceServicer):
    def __init__(self):
        self.model = os.getenv("LLM_MODEL", "gpt-4o-mini")

    async def Generate(self, request, context):
        yield ppt_pb2.PptProgress(status="generating_outline", progress_pct=5,
                                   message="Generating outline...")
        outline = await self._llm_outline(request.topic, request.slide_count)
        prs = Presentation()
        for i, title in enumerate(outline):
            yield ppt_pb2.PptProgress(
                status=f"filling_slide_{i+1}",
                progress_pct=int(30 + 60 * i / len(outline)),
                message=f"Slide {i+1}: {title}")
            content = await self._llm_content(title, request.topic)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = content
        path = f"/tmp/{request.topic.replace(' ', '_')}.pptx"
        prs.save(path)
        yield ppt_pb2.PptProgress(status="done", progress_pct=100,
                                   message="Done!", result_url=path)

    async def _llm_outline(self, topic, count):
        r = await litellm.acompletion(
            model=self.model,
            messages=[{"role": "user",
                       "content": f"List {count or 5} slide titles for '{topic}', comma-separated."}]
        )
        return [s.strip() for s in r.choices[0].message.content.split(",") if s.strip()]

    async def _llm_content(self, title, topic):
        r = await litellm.acompletion(
            model=self.model,
            messages=[{"role": "user",
                       "content": f"Write 3 bullet point content for slide '{title}' about {topic}."}]
        )
        return r.choices[0].message.content

async def main():
    port = int(os.getenv("GRPC_PORT", "50053"))
    server = grpc_server()
    server.add_insecure_port(f"0.0.0.0:{port}")
    ppt_pb2_grpc.add_PptServiceServicer_to_server(PptService(), server)
    print(f"[ai-ppt] starting on port {port}")
    await server.start()
    await server.wait_for_termination()

if __name__ == "__main__":
    asyncio.run(main())
